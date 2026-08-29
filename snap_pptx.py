# -*- coding: utf-8 -*-
"""snap_pptx — ajoute une slide « SNAP » (image + métadonnées KLV) à un deck PowerPoint à partir d'un template
(port serveur, python-pptx, du module COM `powerpoint/mod.rs` du viewer Tauri : mêmes placeholders, même
placement de l'image à la place de la forme « ImageSnaps », image envoyée en arrière-plan).

Template : la slide 1 est le modèle ; la forme dont le nom ou le texte contient « imagesnaps » reçoit l'image ;
les textes `{…}` sont substitués (paragraphe par paragraphe : PowerPoint découpe les placeholders en plusieurs runs).

Placeholders reconnus (insensibles à la casse sur le nom, arguments libres entre parenthèses) :
  {Precision Time Stamp:Date(FMT)}   FMT ∈ ddHHmm · HHmm · HHmmss · MMM · yyyy · yyMMdd · yyyyMMdd · dd · ISO
  {FrameCenter:Geo(MGRS)}            « 31T BL 84414 52901 »      {FrameCenter:Geo(DD)}  « 45.596198 0.235836 »
  {Sensor True Altitude(…)}          « 19672 ft »                {Frame Center Elevation(…)}  « 312 ft »
  {Platform Heading Angle(…)}        « 169.8° »                  {Platform Designation}  indicatif / plateforme
  {Mission}  {Description}  {Callsign}  {Sensor}  {HFOV}  {Slant Range}  {Sensor Lat}  {Sensor Lon}
  _DESCRIPTION (suffixe littéral de la ligne de nommage) → « _<description> » si une description est fournie.
"""
import copy
import os
import re
import threading
import time

from pptx import Presentation
from pptx.util import Emu

_LOCK = threading.Lock()
_PH = re.compile(r"\{([A-Za-z][A-Za-z :]*?)(?::Date)?\(([^)]*)\)\}|\{([A-Za-z][A-Za-z ]*?)\}")
_MONTHS = ["JAN", "FEB", "MAR", "APR", "MAY", "JUN", "JUL", "AUG", "SEP", "OCT", "NOV", "DEC"]


def fmt_date(ts, fmt):
    """Horodatage epoch (UTC) selon les formats du viewer (ddHHmm, HHmm, MMM, yyyy, yyMMdd…) ; ISO sinon."""
    if ts is None:
        return "N/A"
    t = time.gmtime(ts)
    f = (fmt or "").strip()
    table = {"ddHHmm": "%d%H%M", "HHmm": "%H%M", "HHmmss": "%H%M%S", "yyyy": "%Y", "yy": "%y", "yyMMdd": "%y%m%d",
             "yyyyMMdd": "%Y%m%d", "dd": "%d", "MM": "%m", "ddMMyyyy": "%d%m%Y", "HH": "%H", "mm": "%M"}
    if f == "MMM":
        return _MONTHS[t.tm_mon - 1]
    if f in table:
        return time.strftime(table[f], t)
    return time.strftime("%Y-%m-%dT%H:%M:%SZ", t)


def fmt_mgrs(s):
    """'31TBL8441452901' → '31T BL 84414 52901' (zone+bande, carré, easting, northing)."""
    if not s:
        return "N/A"
    m = re.match(r"^(\d{1,2}[A-Z])([A-Z]{2})(\d+)$", s)
    if not m:
        return s
    zb, sq, digits = m.groups(); h = len(digits) // 2
    return "%s %s %s %s" % (zb, sq, digits[:h], digits[h:])


def _num(v, fmt, unit=""):
    return "N/A" if v is None else (fmt % v) + unit


def render_value(name, arg, v):
    """Valeur d'un placeholder à partir du dict `v` (ts, mgrs, alt_ft, fc_elev_ft, hdg, mission, description, callsign…)."""
    n = re.sub(r"\s+", " ", name).strip().lower()
    if n == "precision time stamp":
        return fmt_date(v.get("ts"), arg)
    if n == "framecenter:geo" or n == "frame center:geo" or n == "framecenter":
        if (arg or "").strip().upper() == "DD":
            return "N/A" if v.get("fc_lat") is None else "%.6f %.6f" % (v["fc_lat"], v["fc_lon"])
        return fmt_mgrs(v.get("mgrs"))
    if n == "sensor true altitude":
        return _num(v.get("alt_ft"), "%.0f", " ft")
    if n == "frame center elevation":
        return _num(v.get("fc_elev_ft"), "%.0f", " ft")
    if n == "platform heading angle":
        return _num(v.get("hdg"), "%.1f", "°")
    if n == "platform designation":
        return v.get("callsign") or v.get("platform") or v.get("mission") or "N/A"
    if n == "mission":
        return v.get("mission") or "N/A"
    if n == "description":
        return v.get("description") or ""
    if n == "callsign":
        return v.get("callsign") or "N/A"
    if n == "sensor":
        return v.get("sensor") or "N/A"
    if n == "hfov":
        return _num(v.get("hfov"), "%.2f", "°")
    if n == "slant range":
        return _num(v.get("slant_m"), "%.0f", " m")
    if n == "sensor lat":
        return _num(v.get("lat"), "%.6f")
    if n == "sensor lon":
        return _num(v.get("lon"), "%.6f")
    return None                                            # inconnu : laissé tel quel


def substitute(text, v):
    def rep(m):
        name, arg, bare = m.group(1), m.group(2), m.group(3)
        r = render_value(name if name else bare, arg, v)
        return m.group(0) if r is None else r
    out = _PH.sub(rep, text)
    if v.get("description") and "_DESCRIPTION" in out:
        out = out.replace("_DESCRIPTION", "_" + re.sub(r"[^\w\-]+", "_", v["description"]).strip("_"))
    return out


def _fill_text_frame(tf, v):
    for para in tf.paragraphs:
        runs = para.runs
        if not runs:
            continue
        full = "".join(r.text for r in runs)
        new = substitute(full, v)
        if new != full:
            runs[0].text = new
            for r in runs[1:]:
                r.text = ""


def _walk_shapes(shapes):
    for sh in shapes:
        yield sh
        if sh.shape_type == 6:                                # groupe
            for s in _walk_shapes(sh.shapes):
                yield s


def _clone_slide(prs, src_slide):
    """Nouvelle slide = copie de la slide `src_slide` (formes, relations images) dans la présentation `prs`."""
    layout = None
    for lo in prs.slide_layouts:
        if lo.name == src_slide.slide_layout.name:
            layout = lo; break
    layout = layout or prs.slide_layouts[len(prs.slide_layouts) - 1]
    dst = prs.slides.add_slide(layout)
    for ph in list(dst.shapes):                               # placeholders hérités de la disposition
        ph._element.getparent().remove(ph._element)
    rid_map = {}
    for rel in src_slide.part.rels.values():
        if rel.reltype.endswith("/image") or rel.reltype.endswith("/media") or rel.reltype.endswith("/hyperlink"):
            if rel.is_external:
                rid_map[rel.rId] = dst.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
            else:
                rid_map[rel.rId] = dst.part.relate_to(rel.target_part, rel.reltype)
    for el in src_slide.shapes._spTree.iterchildren():
        tag = el.tag.split("}")[-1]
        if tag in ("nvGrpSpPr", "grpSpPr"):
            continue
        new = copy.deepcopy(el)
        for node in new.iter():
            for attr, val in list(node.attrib.items()):
                if attr.endswith("}embed") or attr.endswith("}link") or attr.endswith("}id"):
                    if val in rid_map:
                        node.set(attr, rid_map[val])
        dst.shapes._spTree.append(new)
    return dst


def _place_image(slide, image_path):
    """Remplace la forme « imagesnaps » par l'image (même cadre, ratio conservé, centrée) envoyée en arrière-plan."""
    target = None
    for sh in _walk_shapes(slide.shapes):
        name = (sh.name or "").lower()
        txt = sh.text_frame.text.lower() if sh.has_text_frame else ""
        if "imagesnaps" in name or "imagesnaps" in txt:
            target = sh; break
    if target is None:
        left, top, width, height = Emu(624881), Emu(725499), Emu(10881236), Emu(6120000)
    else:
        left, top, width, height = target.left, target.top, target.width, target.height
        target._element.getparent().remove(target._element)
    pic = slide.shapes.add_picture(image_path, left, top)
    iw, ih = pic.width, pic.height                             # taille native → ajustée dans le cadre, ratio conservé
    scale = min(width / iw, height / ih)
    pic.width, pic.height = int(iw * scale), int(ih * scale)
    pic.left, pic.top = int(left + (width - pic.width) / 2), int(top + (height - pic.height) / 2)
    tree = slide.shapes._spTree
    tree.remove(pic._element); tree.insert(2, pic._element)    # arrière-plan (après nvGrpSpPr, grpSpPr)
    return pic


def append_capture(deck_path, template_path, image_path, values):
    """Ajoute une slide SNAP au deck (créé depuis le template s'il n'existe pas) ; renvoie le nombre de slides."""
    if not os.path.isfile(template_path):
        raise FileNotFoundError("template PowerPoint introuvable : %s" % template_path)
    with _LOCK:
        tpl = Presentation(template_path)
        if os.path.isfile(deck_path):
            prs = Presentation(deck_path)
        else:
            prs = Presentation(template_path)                   # même thème / dispositions que le template
            for sld in list(prs.slides._sldIdLst):              # deck vide au départ (la slide modèle reste dans le template)
                rid = sld.rId; prs.part.drop_rel(rid); prs.slides._sldIdLst.remove(sld)
        slide = _clone_slide(prs, tpl.slides[0])
        for sh in _walk_shapes(slide.shapes):
            if sh.has_text_frame:
                _fill_text_frame(sh.text_frame, values)
            if sh.shape_type == 19:                             # tableau
                for row in sh.table.rows:
                    for cell in row.cells:
                        _fill_text_frame(cell.text_frame, values)
        _place_image(slide, image_path)
        tmp = deck_path + ".part"
        prs.save(tmp); os.replace(tmp, deck_path)
        return len(prs.slides)


def slide_from_template(template_path, image_path, values, out_path):
    """Présentation d'UNE slide (pour l'agent poste : insertion dans le PowerPoint ouvert)."""
    with _LOCK:
        prs = Presentation(template_path)
        slide = prs.slides[0]
        for sh in _walk_shapes(slide.shapes):
            if sh.has_text_frame:
                _fill_text_frame(sh.text_frame, values)
        _place_image(slide, image_path)
        prs.save(out_path)
    return out_path
